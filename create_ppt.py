#!/usr/bin/env python3
"""创建A股早盘总结PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── 颜色主题 ───
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
MEDIUM_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF4)
BLACK = RGBColor(0x2C, 0x3E, 0x50)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    from pptx.enum.dml import MSO_THEME_COLOR
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_textbox(slide, left, top, width, height, items, font_size=16, color=BLACK, bullet_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = Pt(6)
        p.level = 0
    return txBox

# ════════════════════════════════════════
# 第1页：封面
# ════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK_BLUE)

# 装饰条
add_shape(slide1, Inches(0), Inches(3.0), Inches(13.333), Inches(0.08), MEDIUM_BLUE)

# 主标题
add_textbox(slide1, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
            "A 股 早 盘 速 览", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# 副标题
add_textbox(slide1, Inches(1.5), Inches(3.3), Inches(10), Inches(0.8),
            "2026年6月8日（星期一）· 半日行情总结", font_size=24, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# 底部信息
add_textbox(slide1, Inches(1.5), Inches(5.5), Inches(10), Inches(0.6),
            "数据来源：新浪财经 / 东方财富  |  制图：SuperCodex", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════
# 第2页：核心数据概览
# ════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

# 顶部标题栏
add_shape(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK_BLUE)
add_textbox(slide2, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
            "📊 三大指数早盘表现", font_size=32, bold=True, color=WHITE)

# 三大指数卡片
card_data = [
    ("上证指数", "3,9xx.xx", "-1.26%", RED),
    ("深证成指", "1x,xxx.xx", "-2.0%+", RED),
    ("创业板指", "x,xxx.xx", "-2.83%", RED),
]

for i, (name, value, change, color) in enumerate(card_data):
    left = Inches(1.0 + i * 4.0)
    top = Inches(1.5)
    card = add_shape(slide2, left, top, Inches(3.5), Inches(2.8), LIGHT_GRAY)
    
    # 指数名称
    add_textbox(slide2, left + Inches(0.3), top + Inches(0.3), Inches(2.9), Inches(0.5),
                name, font_size=22, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    
    # 涨跌幅
    add_textbox(slide2, left + Inches(0.3), top + Inches(0.9), Inches(2.9), Inches(0.7),
                change, font_size=36, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    
    # 描述
    desc = "半日跌幅" if i == 0 else ("午盘跌超2%" if i == 1 else "午盘跌2.83%")
    add_textbox(slide2, left + Inches(0.3), top + Inches(1.7), Inches(2.9), Inches(0.5),
                desc, font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# 市场整体数据
add_shape(slide2, Inches(1.0), Inches(4.8), Inches(11.3), Inches(2.2), LIGHT_BLUE)

add_textbox(slide2, Inches(1.3), Inches(4.9), Inches(11), Inches(0.5),
            "📋 市场整体概况", font_size=22, bold=True, color=DARK_BLUE)

bullets = [
    "▸ 全市场超 4,500 只个股下跌，市场情绪整体偏弱",
    "▸ 受韩国 KOSPI 暴跌逾 8% 触发熔断影响，亚太市场承压",
    "▸ 沪深两市半日成交维持缩量，资金观望情绪浓厚",
    "▸ 外资中长期仍看好A股，5月以来155股获海外机构调研",
]
add_bullet_textbox(slide2, Inches(1.3), Inches(5.4), Inches(11), Inches(1.5),
                   bullets, font_size=16, color=BLACK)

# ════════════════════════════════════════
# 第3页：早盘走势回顾
# ════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

# 顶部标题栏
add_shape(slide3, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK_BLUE)
add_textbox(slide3, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
            "⏱ 早盘时间线回顾", font_size=32, bold=True, color=WHITE)

# 时间线
timeline_items = [
    ("09:25", "三大指数集体低开", "半导体板块跌幅居前，市场情绪谨慎"),
    ("09:30-10:00", "探底后小幅回升", "创业板指一度翻红，机器人概念率先走强"),
    ("10:00-10:36", "再度回落，跌幅扩大", "有色、半导体等板块集体下挫，拖累指数"),
    ("10:36-11:00", "创指跌幅收窄", "工程机械板块短暂拉升，但持续性不足"),
    ("11:00-11:30", "午盘加速下跌", "深成指、创业板指均跌超2%，超4500股下跌"),
]

for i, (time, title, desc) in enumerate(timeline_items):
    y = Inches(1.3) + Inches(i * 1.1)
    
    # 时间标签
    time_box = add_shape(slide3, Inches(0.8), y, Inches(1.8), Inches(0.8), DARK_BLUE)
    add_textbox(slide3, Inches(0.8), y + Inches(0.1), Inches(1.8), Inches(0.6),
                time, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # 连接线
    if i < len(timeline_items) - 1:
        line = add_shape(slide3, Inches(1.65), y + Inches(0.8), Inches(0.08), Inches(0.4), MEDIUM_BLUE)
    
    # 标题
    add_textbox(slide3, Inches(3.0), y + Inches(0.0), Inches(4), Inches(0.4),
                title, font_size=18, bold=True, color=DARK_BLUE)
    
    # 描述
    add_textbox(slide3, Inches(3.0), y + Inches(0.4), Inches(9), Inches(0.4),
                desc, font_size=14, color=GRAY)

# ════════════════════════════════════════
# 第4页：板块表现
# ════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)

# 顶部标题栏
add_shape(slide4, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK_BLUE)
add_textbox(slide4, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
            "🔥 板块热点分化", font_size=32, bold=True, color=WHITE)

# 领涨板块 - 左侧
add_shape(slide4, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5), RGBColor(0xE8, 0xF8, 0xF5))
add_textbox(slide4, Inches(1.1), Inches(1.5), Inches(5), Inches(0.5),
            "✅ 领涨板块", font_size=24, bold=True, color=GREEN)

gainers = [
    "🤖  机器人（减速器）概念 — 逆市爆发",
    "     • 凡拓数创涨停，津膜科技涨幅居前",
    "     • 机器人ETF易方达涨3.62%",
    "     • 产业迎密集催化：4天4笔融资",
    "",
    "⛏  煤炭开采加工 — 短线拉升",
    "     • 复产慢、供需偏紧推动板块上行",
    "     • 大有能源5连板（昨日）",
    "",
    "🏗  工程机械 — 盘中异动",
    "     • 5月挖机增长超预期",
    "     • 出海保持高景气度",
]
add_bullet_textbox(slide4, Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.5),
                   gainers, font_size=15, color=BLACK)

# 领跌板块 - 右侧
add_shape(slide4, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.5), RGBColor(0xFD, 0xED, 0xEC))
add_textbox(slide4, Inches(7.1), Inches(1.5), Inches(5), Inches(0.5),
            "❌ 领跌板块", font_size=24, bold=True, color=RED)

losers = [
    "🔴  有色金属 — 集体下挫",
    "     • 板块资金流出明显",
    "     • 受全球避险情绪拖累",
    "",
    "🔴  半导体 — 板块调整",
    "     • 科技赛道拥挤度引发热议",
    "     • 算力板块同步回调",
    "",
    "🔴  MLCC / 商业航天",
    "     • MLCC概念领跌",
    "     • 商业航天概念调整",
    "",
    "🔴  算力 / AI概念",
    "     • 科技板块大幅波动",
    "     • 资金获利了结压力",
]
add_bullet_textbox(slide4, Inches(7.1), Inches(2.1), Inches(5.0), Inches(4.5),
                   losers, font_size=15, color=BLACK)

# ════════════════════════════════════════
# 第5页：背景与展望
# ════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)

# 顶部标题栏
add_shape(slide5, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK_BLUE)
add_textbox(slide5, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
            "🌍 市场背景与午后关注", font_size=32, bold=True, color=WHITE)

# 左侧：国际市场背景
add_shape(slide5, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.5), LIGHT_GRAY)
add_textbox(slide5, Inches(1.1), Inches(1.5), Inches(5.2), Inches(0.5),
            "🌐 国际市场扰动", font_size=22, bold=True, color=DARK_BLUE)

bg_items = [
    "▸ 韩国KOSPI指数暴跌逾8%，触发熔断",
    '     —— 亚太市场"黑色星期一"情绪蔓延',
    "",
    "▸ 伊朗空袭以色列事件升级",
    "     —— 特朗普紧急发声，地缘风险升温",
    "",
    '▸ 美股上周五"黑色星期五"',
    "     —— 全球避险情绪传导至A股早盘",
    "",
    "▸ 外资动向：中长期仍看好A股",
    "     —— QFII青睐电子等新质生产力赛道",
    "     —— 5月以来155股获海外机构调研",
]
add_bullet_textbox(slide5, Inches(1.1), Inches(2.1), Inches(5.2), Inches(4.5),
                   bg_items, font_size=15, color=BLACK)

# 右侧：午后关注
add_shape(slide5, Inches(7.2), Inches(1.4), Inches(5.3), Inches(5.5), RGBColor(0xFD, 0xF2, 0xE9))
add_textbox(slide5, Inches(7.5), Inches(1.5), Inches(4.8), Inches(0.5),
            "🔍 午后关注要点", font_size=22, bold=True, color=ORANGE)

focus_items = [
    "① 护盘力量是否出现",
    "    金融、权重股能否企稳",
    "",
    "② 机器人主线能否延续",
    "    减速器概念持续性是关键",
    "",
    "③ 北向资金流向变化",
    "    外资午后是否会逆势加仓",
    "",
    "④ 外围市场企稳信号",
    "    韩股、美股期货走势参考",
    "",
    "⑤ 政策面消息",
    "    国常会研究未来产业发展",
    "    算电协同+迎峰度夏催化",
]
add_bullet_textbox(slide5, Inches(7.5), Inches(2.1), Inches(4.8), Inches(4.5),
                   focus_items, font_size=15, color=BLACK)

# ════════════════════════════════════════
# 第6页：结语 / 风险提示
# ════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, DARK_BLUE)

add_shape(slide6, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), MEDIUM_BLUE)

add_textbox(slide6, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
            "📌 小结", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

summary_items = [
    "今日A股三大指数早盘集体低开低走，受外围市场暴跌及地缘风险拖累，",
    "超4500只个股下跌。机器人（减速器）概念逆市爆发成为唯一亮点。",
    "午后关注护盘力度、热点持续性及外围企稳信号。",
]
add_bullet_textbox(slide6, Inches(1.5), Inches(3.6), Inches(10), Inches(1.5),
                   summary_items, font_size=20, color=LIGHT_BLUE)

# 风险提示
add_textbox(slide6, Inches(1.5), Inches(5.5), Inches(10), Inches(0.8),
            "⚠️ 风险提示：以上内容仅供参考，不构成投资建议。市场有风险，投资需谨慎。",
            font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# ─── 保存 ───
output_path = "/Users/a1021500689/Documents/SuperCodex/A股早盘速览_20260608.pptx"
prs.save(output_path)
print(f"✅ PPT 已保存至: {output_path}")
